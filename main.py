import asyncio
import streamlit as st
import os
import tempfile
from pptx import Presentation
from googletrans import Translator
import base64
import subprocess

from utils.init import initialize
from utils.counter import initialize_user_count, increment_user_count, decrement_user_count, get_user_count
from utils.TelegramSender import TelegramSender
# from utils.tools import get_image_url

def convert_to_pptx(input_file):
    file_name, file_extension = os.path.splitext(input_file)
    output_file = file_name + ".pptx"
    
    if file_extension.lower() in ['.ppt', '.odp', '.otp']:
        # המרה ל-pptx באמצעות LibreOffice
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pptx', input_file, '--outdir', os.path.dirname(output_file)])
        return output_file
    elif file_extension.lower() == '.pptx':
        return input_file
    else:
        # עבור פורמטים לא נתמכים, נזרוק שגיאה
        raise ValueError(f"פורמט קובץ לא נתמך: {file_extension}")

def translate_text(text, dest='he'):
    translator = Translator()
    try:
        return translator.translate(text, dest=dest).text
    except Exception as e:
        st.error(f"שגיאת תרגום: {e}")
        return text

def translate_slide(slide):
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    if original_text.strip():
                        translated_text = translate_text(original_text)
                        run.text = translated_text

def translate_pptx(input_file):
    prs = Presentation(input_file)
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    for i, slide in enumerate(prs.slides):
        status_text.text(f"מתרגם שקופית {i+1}/{len(prs.slides)}")
        translate_slide(slide)
        progress_bar.progress((i + 1) / len(prs.slides))
    
    # status_text.text("שומר מצגת מתורגמת...")
    
    # שמירה לקובץ זמני
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)
    temp_file.close()
    
    return temp_file.name

def get_binary_file_downloader_html(bin_file, file_label='File'):
    with open(bin_file, 'rb') as f:
        data = f.read()
    bin_str = base64.b64encode(data).decode()
    href = f'data:application/octet-stream;base64,{bin_str}'
    custom_css = f"""
        <style>
            .download-button {{
                width: 100%;
                margin: 0.5rem 0;
                background-color: #ff0000;  /* אדום */
                color: white !important;  /* לבן עם !important כדי לוודא שהצבע יחול */
                font-weight: bold;
                border: 2px solid #ff0000;
                border-radius: 10px;
                padding: 12px 24px;
                text-align: center;
                text-decoration: none;
                display: inline-block;
                font-size: 24px;
                cursor: pointer;
                transition: background-color 0.3s, border-color 0.3s;
                
            }}
            .download-button:hover {{
                background-color: #cc0000;  /* אדום כהה יותר */
                border-color: #cc0000;
                color: white !important;  /* לוודא שהצבע נשאר לבן גם ב-hover */
            }}
            @media screen and (max-width: 768px) {{
                .download-button {{
                    padding: 10px 20px;
                    font-size: 16px;
                }}
            }}
            @media screen and (max-width: 480px) {{
                .download-button {{
                    padding: 8px 16px;
                    font-size: 14px;
                }}
            }}
        </style>
    """
    download_button = f'<a href="{href}" download="{os.path.basename(bin_file)}" class="download-button">הורד {file_label}</a>'
    return f'{custom_css}{download_button}'

# Initialize TelegramSender
if 'telegram_sender' not in st.session_state:
    st.session_state.telegram_sender = TelegramSender()

# Increment user count if this is a new session
if 'counted' not in st.session_state:
    st.session_state.counted = True
    increment_user_count()

# Initialize user count
initialize_user_count()

def main():
    header_content, image_path, footer_content = initialize()

    st.markdown(f"<h2 style='text-align: center; color: #FF6347;'>{header_content}</h2>", unsafe_allow_html=True)

    if image_path:
        st.image(image_path, use_column_width=True)

    if 'translated_file' not in st.session_state:
        st.session_state.translated_file = None

    uploaded_file = st.file_uploader("בחרו קובץ PowerPoint", type=["pptx", "ppt", "odp", "otp"])

    if uploaded_file is not None:
        # st.write("הקובץ הועלה בהצלחה!")
        
        if st.button("תרגם את המצגת לעברית"):
            with st.spinner("מתרגם לעברית..."):
                # שמירת הקובץ שהועלה באופן זמני
                temp_input = tempfile.NamedTemporaryFile(delete=False, suffix="." + uploaded_file.name.split('.')[-1])
                temp_input.write(uploaded_file.read())
                temp_input.close()

                try:
                    # המרה ל-PPTX אם נדרש
                    pptx_file = convert_to_pptx(temp_input.name)
                    
                    # תרגום הקובץ
                    output_file = translate_pptx(pptx_file)

                    # שמירת נתיב קובץ הפלט ב-session state
                    st.session_state.translated_file = output_file

                    # שליחה לטלגרם
                    asyncio.run(send_telegram_message_and_file(
                        f"New PowerPoint file translated: {uploaded_file.name}",
                        output_file
                    ))                    

                    # ניקוי קבצים זמניים
                    os.unlink(temp_input.name)
                    if pptx_file != temp_input.name:
                        os.unlink(pptx_file)

                    st.success("התרגום הושלם ניתן להוריד!")
                except Exception as e:
                    st.error(f"אירעה שגיאה: {e}")

        # הצגת קישור להורדה אם קובץ תורגם
        if st.session_state.translated_file:            
            st.markdown(get_binary_file_downloader_html(st.session_state.translated_file, 'מצגת מתורגמת'), unsafe_allow_html=True)
    
    user_count = get_user_count(formatted=True)
    footer_with_count = f"{footer_content}\n\n<p class='user-count' style='color: #4B0082;'>סה\"כ משתמשים: {user_count}</p>"
    st.markdown(footer_with_count, unsafe_allow_html=True)

async def send_telegram_message_and_file(message, file_path):
    sender = st.session_state.telegram_sender
    try:
        await sender.send_document(file_path, message)
    finally:
        await sender.close_session()

if __name__ == "__main__":
    main()
